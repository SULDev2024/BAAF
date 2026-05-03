"""
Build Topic 15 academic deck (16:9) + figures from project data.
Run: python presentation/build_topic15_deck.py
Outputs: presentation/Topic15_Bioinformatics_Project.pptx + presentation/assets/*.png
"""
from __future__ import annotations

import os
import sys

import matplotlib.pyplot as plt
import numpy as np
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
from scipy.cluster.hierarchy import dendrogram, linkage
from scipy.spatial.distance import squareform

ROOT = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
ASSETS = os.path.join(os.path.dirname(__file__), "assets")
OUT_PPTX = os.path.join(os.path.dirname(__file__), "Topic15_Bioinformatics_Project.pptx")

TEAL = RGBColor(13, 148, 136)
SLATE = RGBColor(15, 23, 42)
MUTED = RGBColor(71, 85, 105)
WHITE = RGBColor(255, 255, 255)

# Worked example from report (taxa A,B,C,D)
D_MATRIX = np.array(
    [
        [0, 5, 9, 9],
        [5, 0, 10, 10],
        [9, 10, 0, 8],
        [9, 10, 8, 0],
    ],
    dtype=float,
)
TAXA = ["A", "B", "C", "D"]


def ensure_assets() -> tuple[str, str]:
    os.makedirs(ASSETS, exist_ok=True)
    heat_path = os.path.join(ASSETS, "distance_matrix_heatmap.png")
    tree_path = os.path.join(ASSETS, "upgma_dendrogram.png")

    plt.rcParams.update(
        {
            "font.family": "sans-serif",
            "font.sans-serif": ["Segoe UI", "Arial", "DejaVu Sans"],
            "axes.titlesize": 14,
            "figure.facecolor": "white",
        }
    )

    # Heatmap
    fig, ax = plt.subplots(figsize=(6, 5.2))
    im = ax.imshow(D_MATRIX, cmap="YlGnBu", vmin=0, vmax=12)
    ax.set_xticks(range(4), TAXA)
    ax.set_yticks(range(4), TAXA)
    for i in range(4):
        for j in range(4):
            ax.text(j, i, int(D_MATRIX[i, j]), ha="center", va="center", color="white" if D_MATRIX[i, j] > 6 else "#0f172a", fontsize=12)
    ax.set_title("Aggregate distance matrix (example dataset)")
    plt.colorbar(im, ax=ax, fraction=0.046, pad=0.04, label="Distance")
    plt.tight_layout()
    fig.savefig(heat_path, dpi=200, bbox_inches="tight")
    plt.close(fig)

    # UPGMA / average linkage dendrogram (matches project implementation concept)
    condensed = squareform(D_MATRIX, checks=False)
    Z = linkage(condensed, method="average")
    fig, ax = plt.subplots(figsize=(7.5, 4.2))
    dendrogram(Z, labels=TAXA, ax=ax, color_threshold=0, above_threshold_color="#0d9488")
    ax.set_title("UPGMA-style hierarchical clustering (average linkage)")
    ax.set_ylabel("Merge height")
    ax.spines["top"].set_visible(False)
    ax.spines["right"].set_visible(False)
    plt.tight_layout()
    fig.savefig(tree_path, dpi=200, bbox_inches="tight")
    plt.close(fig)

    return heat_path, tree_path


def _blank_slide(prs: Presentation):
    return prs.slides.add_slide(prs.slide_layouts[6])


def _header_bar(slide, title: str, slide_width):
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, slide_width, Inches(1.05)
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = TEAL
    bar.line.fill.background()
    box = slide.shapes.add_textbox(Inches(0.45), Inches(0.22), slide_width - Inches(0.6), Inches(0.7))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.LEFT


def _body_box(slide, top_in: float, left_in: float = 0.45, width_in: float = 12.4):
    box = slide.shapes.add_textbox(Inches(left_in), Inches(top_in), Inches(width_in), Inches(5.8))
    tf = box.text_frame
    tf.word_wrap = True
    return tf


def _add_bullets(tf, lines: list[str], size=18):
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.font.size = Pt(size)
        p.font.color.rgb = SLATE
        p.space_after = Pt(6)
        p.level = 0


def build_prs(heat_path: str, tree_path: str) -> Presentation:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # 1 Title
    s = _blank_slide(prs)
    title = s.shapes.add_textbox(Inches(0.7), Inches(2.0), Inches(12), Inches(1.3))
    tp = title.text_frame.paragraphs[0]
    tp.text = "Distance-Based & Character-Based\nPhylogenetic Reconstruction"
    tp.font.size = Pt(40)
    tp.font.bold = True
    tp.font.color.rgb = SLATE
    tp.alignment = PP_ALIGN.LEFT
    sub = s.shapes.add_textbox(Inches(0.7), Inches(3.55), Inches(12), Inches(1.5))
    sp = sub.text_frame.paragraphs[0]
    sp.text = (
        "Topic 15 · Bioinformatics (2025/2026)\n"
        "UPGMA · Neighbor Joining · Fitch parsimony · Django web platform"
    )
    sp.font.size = Pt(20)
    sp.font.color.rgb = MUTED
    sp.alignment = PP_ALIGN.LEFT
    foot = s.shapes.add_textbox(Inches(0.7), Inches(6.35), Inches(12), Inches(0.6))
    fp = foot.text_frame.paragraphs[0]
    fp.text = "University project · Information Technologies & AI · Computer Engineering"
    fp.font.size = Pt(14)
    fp.font.color.rgb = MUTED

    slides_spec: list[tuple[str, list[str]]] = [
        (
            "What this project covers",
            [
                "Phylogenetic trees from aggregate distance matrices (UPGMA, NJ)",
                "Character-based scoring with Fitch parsimony on four taxa",
                "Minimum vs maximum frugality across three candidate topologies",
                "Two deliverables: educational Python script + Django web application",
            ],
        ),
        (
            "Why it matters",
            [
                "Evolutionary relationships drive genomics, epidemiology, and comparative biology",
                "Distances compress signal; characters preserve site-level information",
                "Method choice depends on assumptions (e.g., ultrametricity / molecular clock)",
            ],
        ),
        (
            "UPGMA — hierarchical clustering",
            [
                "Agglomerative: repeatedly merge the closest pair of clusters",
                "Update distances by arithmetic mean across all cross-cluster pairs",
                "Merge height guides branch lengths; outputs rooted-style Newick",
                "Best interpreted when distances are close to ultrametric",
            ],
        ),
        (
            "Neighbor Joining (NJ)",
            [
                "Still distance-based, but relaxes strict ultrametric expectations",
                "Uses row sums r(i) and Q(i,j) to select the next join (not always smallest d)",
                "Computes limb lengths and updates distances to the new internal node",
                "Useful when rate variation breaks clock-like structure",
            ],
        ),
        (
            "Fitch parsimony & frugality",
            [
                "Four taxa → three unrooted binary splits evaluated exhaustively",
                "Per site: intersect child state sets when possible; else union (+1 change)",
                "Minimum frugality = smallest total changes (parsimony optimum)",
                "Maximum frugality = contrast topology among the same three candidates (Topic 15)",
            ],
        ),
        (
            "Worked distance example (A–D)",
            [
                "Symmetric matrix, zero diagonal, non-negative entries (validated in code)",
                "Merge order: (A,B) then (C,D), then join super-clusters — see dendrogram figure",
                "Representative Newick (UPGMA):\n((A:2.50,B:2.50):2.25,(C:4.00,D:4.00):0.75);",
            ],
        ),
        (
            "Character example (four sequences)",
            [
                "Format: Taxon:SEQUENCE with equal length and alphabetic symbols",
                "A:ATGC · B:ATGT · C:CTGC · D:CTGT",
                "Scores computed per topology; min and max frugality reported for comparison",
            ],
        ),
        (
            "Software artifacts",
            [
                "simple-algorithm/topic15_simple.py — compact, defense-friendly trace of logic",
                "phylo-webapp/ — Django services, persistence, authentication, interactive UI",
                "Shared concepts: same UPGMA / NJ / Fitch ideas in script and web layers",
            ],
        ),
        (
            "Reliability & outputs",
            [
                "Matrix checks: square, symmetric, zero diagonal, non-negative",
                "Character checks: four taxa, equal length, strict parsing rules",
                "Outputs: step logs, Newick strings, SVG tree views (e.g., simple-algorithm-output/)",
            ],
        ),
        (
            "Web platform highlights",
            [
                "Run distance-based and character-based analyses from a browser",
                "Store run history with authenticated access",
                "Inspect intermediate and final results suitable for lab demonstration",
            ],
        ),
        (
            "Strengths & limitations",
            [
                "Strengths: deterministic, transparent steps, dual method coverage",
                "Limits: parsimony mode fixed to four taxa; exhaustive search not scalable",
                "Future: file ingest (FASTA/CSV), larger topology search, support metrics",
            ],
        ),
        (
            "Conclusion",
            [
                "Unified treatment of Topic 15: theory → algorithms → software",
                "Clear comparison between distance reduction and character-level inference",
                "Thank you — questions?",
            ],
        ),
    ]

    for title_text, bullets in slides_spec:
        s = _blank_slide(prs)
        _header_bar(s, title_text, prs.slide_width)
        tf = _body_box(s, 1.25)
        _add_bullets(tf, bullets, size=17)

    # Visual slide: matrix + dendrogram
    s = _blank_slide(prs)
    _header_bar(s, "Visual summary · matrix & UPGMA-style tree", prs.slide_width)
    s.shapes.add_picture(heat_path, Inches(0.45), Inches(1.2), width=Inches(5.8))
    s.shapes.add_picture(tree_path, Inches(6.55), Inches(1.15), width=Inches(6.35))

    cap = s.shapes.add_textbox(Inches(0.45), Inches(6.45), Inches(12.4), Inches(0.85))
    cp = cap.text_frame.paragraphs[0]
    cp.text = "Figures generated from the report’s example matrix; dendrogram uses average linkage (UPGMA concept)."
    cp.font.size = Pt(12)
    cp.font.color.rgb = MUTED
    cp.font.italic = True

    # Canva enhancement slide (instructions in speaker notes)
    s = _blank_slide(prs)
    _header_bar(s, "Polish in Canva (recommended)", prs.slide_width)
    tf = _body_box(s, 1.25)
    _add_bullets(
        tf,
        [
            "Import this .pptx via Canva → Create a design → Import → Microsoft PowerPoint",
            "Replace figures with licensed Canva photos/illustrations if you prefer stock art",
            "Suggested element searches: “DNA helix minimal”, “phylogenetic tree line art”, “data matrix abstract”, “genomics gradient”",
            "Keep one accent color (teal) + dark slate text for a cohesive academic look",
        ],
        size=17,
    )
    notes = s.notes_slide.notes_text_frame
    notes.text = (
        "Speaker: Mention that the Canva Developer MCP in Cursor supports app-building docs, "
        "not remote editing of personal Canva decks — this file is built for import into Canva."
    )

    return prs


def main():
    heat_path, tree_path = ensure_assets()
    prs = build_prs(heat_path, tree_path)
    prs.save(OUT_PPTX)
    print(f"Saved: {OUT_PPTX}")


if __name__ == "__main__":
    sys.path.insert(0, ROOT)
    main()
